"""Source ingestion and chunking for carl learn.

Supports: URL, file, directory, HuggingFace dataset, raw text,
          Claude Code skill files (skill://), MCP server definitions (mcp://),
          OpenClaw datasets (claw://).
Chunking strategy adapts to content type:
  - Python/code: module-level boundaries (class/function defs)
  - Markdown: heading boundaries
  - Plain text: paragraph boundaries
Target chunk size: ~500-2000 tokens (estimated as chars / 4).
"""

from __future__ import annotations

import csv
import logging
import os
import re
from enum import Enum
from pathlib import Path
from typing import Any, Optional

from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


class SourceType(str, Enum):
    URL = "url"
    FILE = "file"
    DIRECTORY = "directory"
    HUB_DATASET = "hub_dataset"
    TEXT = "text"
    SKILL = "skill"   # skill://path/to/skill.md
    MCP = "mcp"       # mcp://server-name
    CLAW = "claw"     # claw://dataset-name
    NATURAL = "natural"  # natural language query


class SourceChunk(BaseModel):
    """A chunk of source material."""

    text: str
    source: str = Field(description="Origin path, URL, or identifier")
    chunk_id: int
    modality: str = "text"
    metadata: dict[str, Any] = Field(default_factory=dict)


# Token estimate: ~4 chars per token
_MIN_CHUNK_CHARS = 200
_TARGET_CHUNK_CHARS = 3000
_MAX_CHUNK_CHARS = 8000

# File extensions to ingest from directories
_TEXT_EXTS = {".py", ".md", ".txt", ".json", ".rst", ".yaml", ".yml", ".toml"}
_PDF_EXTS = {".pdf"}
_TABULAR_EXTS = {".csv", ".tsv"}
_OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}
_AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
_VIDEO_EXTS = {".mp4", ".webm", ".mov", ".mkv", ".avi"}
_INGESTABLE_EXTS = _TEXT_EXTS | _PDF_EXTS | _TABULAR_EXTS | _OFFICE_EXTS | _IMAGE_EXTS | _AUDIO_EXTS | _VIDEO_EXTS


class SourceIngester:
    """Reads and chunks source material."""

    def ingest(
        self, source: str, source_type: Optional[SourceType] = None
    ) -> list[SourceChunk]:
        """Auto-detect source type if not provided, then ingest and chunk.

        Args:
            source: Path, URL, HF dataset ID, or raw text.
            source_type: Explicit type override. Auto-detected if None.

        Returns:
            List of SourceChunk with chunk_id starting at 0.

        Raises:
            ValueError: If source cannot be read or type is unrecognized.
            FileNotFoundError: If file/directory does not exist.
        """
        if source_type is None:
            source_type = self._detect_type(source)

        if source_type == SourceType.TEXT:
            return self._chunk_text(source, origin="<text>")
        elif source_type == SourceType.FILE:
            if source.endswith((".zip", ".tar.gz", ".tgz")):
                return self._ingest_archive(source)
            return self._ingest_file(source)
        elif source_type == SourceType.DIRECTORY:
            if source.startswith(("https://github.com/", "http://github.com/")):
                return self._ingest_github(source)
            return self._ingest_directory(source)
        elif source_type == SourceType.URL:
            return self._ingest_url(source)
        elif source_type == SourceType.HUB_DATASET:
            return self._ingest_hub_dataset(source)
        elif source_type == SourceType.SKILL:
            return self._ingest_skill(source)
        elif source_type == SourceType.MCP:
            return self._ingest_mcp(source)
        elif source_type == SourceType.CLAW:
            return self._ingest_claw(source)
        elif source_type == SourceType.NATURAL:
            raise ValueError(
                "Natural language detected. Use interpret_natural() from "
                "carl_studio.learn.planner to get a LearningPlan first."
            )
        else:
            raise ValueError(f"Unknown source type: {source_type!r}")

    # ------------------------------------------------------------------
    # Type detection
    # ------------------------------------------------------------------

    @staticmethod
    def _detect_type(source: str) -> SourceType:
        """Heuristic source type detection."""
        if source.startswith("skill://"):
            return SourceType.SKILL
        if source.startswith("mcp://"):
            return SourceType.MCP
        if source.startswith("claw://"):
            return SourceType.CLAW

        # GitHub repo (must check before generic URL)
        if source.startswith(("https://github.com/", "http://github.com/")):
            return SourceType.DIRECTORY  # will be cloned first

        if source.startswith(("http://", "https://")):
            return SourceType.URL

        # Archive files
        if source.endswith((".zip", ".tar.gz", ".tgz")):
            return SourceType.FILE  # will be extracted first

        path = Path(source)
        if path.is_dir():
            return SourceType.DIRECTORY
        if path.is_file():
            return SourceType.FILE

        # HF dataset pattern: org/dataset or org/dataset-name
        if re.match(r"^[\w.-]+/[\w.-]+$", source) and not path.exists():
            return SourceType.HUB_DATASET

        # Fallback: if it contains newlines or is long, treat as text
        if "\n" in source or len(source) > 200:
            return SourceType.TEXT

        # Natural language: has spaces, long enough to be a sentence
        if " " in source and len(source) > 10:
            return SourceType.NATURAL

        raise ValueError(
            f"Cannot detect source type for {source!r}. "
            "Provide source_type explicitly."
        )

    # ------------------------------------------------------------------
    # Ingestion methods
    # ------------------------------------------------------------------

    def _ingest_file(self, path: str) -> list[SourceChunk]:
        """Read a single file and chunk it, dispatching by extension."""
        p = Path(path)
        if not p.is_file():
            raise FileNotFoundError(f"File not found: {path}")
        return self._dispatch_file(p)

    def _dispatch_file(self, p: Path, start_id: int = 0) -> list[SourceChunk]:
        """Route a file to the right handler based on extension."""
        ext = p.suffix.lower()
        if ext in _PDF_EXTS:
            return self._ingest_pdf(str(p), start_id)
        if ext in _TABULAR_EXTS:
            return self._ingest_tabular(str(p), start_id)
        if ext in _OFFICE_EXTS:
            return self._ingest_office(str(p), start_id)
        if ext in _IMAGE_EXTS:
            return self._ingest_image(str(p), start_id)
        if ext in _AUDIO_EXTS:
            return self._ingest_audio(str(p), start_id)
        if ext in _VIDEO_EXTS:
            return self._ingest_video(str(p), start_id)
        # Default: text
        text = p.read_text(encoding="utf-8", errors="replace")
        if not text.strip():
            return []
        return self._chunk_text(text, origin=str(p), start_id=start_id)

    def _ingest_directory(self, path: str) -> list[SourceChunk]:
        """Glob ingestable files from a directory, chunk each."""
        p = Path(path)
        if not p.is_dir():
            raise FileNotFoundError(f"Directory not found: {path}")

        chunks: list[SourceChunk] = []
        files = sorted(
            f for f in p.rglob("*") if f.is_file() and f.suffix.lower() in _INGESTABLE_EXTS
        )
        if not files:
            raise ValueError(
                f"No ingestable files found in {path}. "
                f"Supported extensions: {sorted(_INGESTABLE_EXTS)}"
            )

        chunk_id = 0
        for f in files:
            try:
                file_chunks = self._dispatch_file(f, start_id=chunk_id)
                chunks.extend(file_chunks)
                chunk_id += len(file_chunks)
            except Exception as exc:
                logger.warning("Skipping %s: %s", f, exc)
                continue
        return chunks

    def _ingest_url(self, url: str) -> list[SourceChunk]:
        """Fetch URL content and chunk it."""
        import urllib.request
        import urllib.error

        try:
            req = urllib.request.Request(
                url, headers={"User-Agent": "carl-studio/0.2.0"}
            )
            with urllib.request.urlopen(req, timeout=30) as resp:
                text = resp.read().decode("utf-8", errors="replace")
        except urllib.error.URLError as exc:
            raise ValueError(f"Failed to fetch URL {url}: {exc}") from exc

        if not text.strip():
            return []

        # Strip HTML tags if content looks like HTML
        if "<html" in text.lower()[:500]:
            text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
            text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()

        return self._chunk_text(text, origin=url)

    def _ingest_hub_dataset(self, dataset_id: str) -> list[SourceChunk]:
        """Load a HuggingFace dataset and extract text fields."""
        try:
            from datasets import load_dataset
        except ImportError as exc:
            raise ImportError(
                "Install 'datasets' to ingest HuggingFace datasets: "
                "pip install datasets"
            ) from exc

        ds = load_dataset(dataset_id, split="train")
        text_fields = [
            col
            for col in ds.column_names
            if col in ("text", "content", "question", "answer", "instruction", "output")
        ]
        if not text_fields:
            raise ValueError(
                f"No text fields found in {dataset_id}. "
                f"Columns: {ds.column_names}"
            )

        chunks: list[SourceChunk] = []
        chunk_id = 0
        for i, row in enumerate(ds):
            parts = [str(row[f]) for f in text_fields if row.get(f)]
            combined = "\n\n".join(parts)
            if len(combined) < _MIN_CHUNK_CHARS:
                # Small rows: accumulate until we hit target
                if chunks and len(chunks[-1].text) < _TARGET_CHUNK_CHARS:
                    chunks[-1] = SourceChunk(
                        text=chunks[-1].text + "\n\n" + combined,
                        source=f"{dataset_id}[{chunks[-1].source.split('[')[1]}",
                        chunk_id=chunks[-1].chunk_id,
                    )
                    continue
            chunk = SourceChunk(
                text=combined, source=f"{dataset_id}[{i}]", chunk_id=chunk_id
            )
            chunks.append(chunk)
            chunk_id += 1
        return chunks

    def _ingest_skill(self, source: str) -> list[SourceChunk]:
        """Ingest a Claude Code skill file. skill://path/to/skill.md"""
        path = source.removeprefix("skill://")
        if not os.path.exists(path):
            raise FileNotFoundError(f"Skill file not found: {path}")
        with open(path, encoding="utf-8", errors="replace") as f:
            content = f.read()
        if not content.strip():
            return []
        return self._chunk_text(content, origin=path)

    def _ingest_mcp(self, source: str) -> list[SourceChunk]:
        """Ingest MCP server tool definitions. mcp://server-name

        Reads MCP config from ~/.claude.json or .mcp.json, extracts tool
        schemas for the named server.
        """
        import json

        server_name = source.removeprefix("mcp://")
        config_paths = [
            Path.home() / ".claude.json",
            Path(".mcp.json"),
            Path(".claude") / "mcp.json",
        ]
        tools_text = f"# MCP Server: {server_name}\n\nServer tools to learn:\n"
        found = False
        for config_path in config_paths:
            if config_path.exists():
                try:
                    with open(config_path, encoding="utf-8") as f:
                        config = json.load(f)
                    servers = config.get("mcpServers", config.get("servers", {}))
                    if server_name in servers:
                        server = servers[server_name]
                        cmd = server.get("command", "?")
                        args = " ".join(server.get("args", []))
                        tools_text += f"\nCommand: {cmd} {args}\n"
                        env_vars = server.get("env", {})
                        if env_vars:
                            tools_text += f"Environment: {', '.join(env_vars.keys())}\n"
                        found = True
                        break
                except (json.JSONDecodeError, OSError):
                    continue
        if not found:
            tools_text += f"\n(No config found for server '{server_name}')\n"
        return [SourceChunk(text=tools_text, source=f"mcp://{server_name}", chunk_id=0)]

    # ------------------------------------------------------------------
    # Multimodal handlers (graceful degradation — skip if deps missing)
    # ------------------------------------------------------------------

    def _ingest_pdf(self, path: str, start_id: int = 0) -> list[SourceChunk]:
        """Extract text from PDF. Tries pymupdf, then pypdf, then skip."""
        text = ""
        try:
            import pymupdf  # noqa: F811

            doc = pymupdf.open(path)
            text = "\n\n".join(page.get_text() for page in doc)
            doc.close()
        except ImportError:
            try:
                from pypdf import PdfReader

                reader = PdfReader(path)
                text = "\n\n".join(
                    page.extract_text() or "" for page in reader.pages
                )
            except ImportError:
                logger.warning(
                    "Skipping PDF %s: install pymupdf or pypdf  "
                    "(pip install pymupdf  or  pip install pypdf)", path,
                )
                return []
        if not text.strip():
            return []
        chunks = self._chunk_text(text, origin=path, start_id=start_id)
        for c in chunks:
            c.modality = "document"
        return chunks

    def _ingest_tabular(self, path: str, start_id: int = 0) -> list[SourceChunk]:
        """CSV/TSV → markdown table representation. Stdlib only."""
        p = Path(path)
        delimiter = "\t" if p.suffix == ".tsv" else ","
        text_parts: list[str] = []
        with open(p, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            rows = list(reader)
        if not rows:
            return []
        header = rows[0]
        text_parts.append("| " + " | ".join(header) + " |")
        text_parts.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            text_parts.append("| " + " | ".join(row) + " |")
        text = "\n".join(text_parts)
        chunks = self._chunk_text(text, origin=path, start_id=start_id)
        for c in chunks:
            c.modality = "tabular"
        return chunks

    def _ingest_office(self, path: str, start_id: int = 0) -> list[SourceChunk]:
        """Extract text from .docx/.xlsx/.pptx. Graceful if deps missing."""
        p = Path(path)
        ext = p.suffix.lower()
        text = ""
        if ext == ".docx":
            try:
                from docx import Document  # type: ignore[import-untyped]

                doc = Document(path)
                text = "\n\n".join(para.text for para in doc.paragraphs if para.text.strip())
            except ImportError:
                logger.warning("Skipping %s: pip install python-docx", path)
                return []
        elif ext == ".xlsx":
            try:
                from openpyxl import load_workbook  # type: ignore[import-untyped]

                wb = load_workbook(path, read_only=True, data_only=True)
                parts: list[str] = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(c) if c is not None else "" for c in row]
                        if any(vals):
                            parts.append(" | ".join(vals))
                text = "\n".join(parts)
                wb.close()
            except ImportError:
                logger.warning("Skipping %s: pip install openpyxl", path)
                return []
        elif ext == ".pptx":
            try:
                from pptx import Presentation  # type: ignore[import-untyped]

                prs = Presentation(path)
                parts_p: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts_p.append(shape.text_frame.text)
                text = "\n\n".join(parts_p)
            except ImportError:
                logger.warning("Skipping %s: pip install python-pptx", path)
                return []
        if not text.strip():
            return []
        chunks = self._chunk_text(text, origin=path, start_id=start_id)
        for c in chunks:
            c.modality = "document"
        return chunks

    def _ingest_image(self, path: str, start_id: int = 0) -> list[SourceChunk]:
        """OCR an image file. Tries pytesseract, then skip."""
        try:
            import pytesseract  # type: ignore[import-untyped]
            from PIL import Image  # type: ignore[import-untyped]

            img = Image.open(path)
            text = pytesseract.image_to_string(img)
            img.close()
        except ImportError:
            logger.warning(
                "Skipping image %s: pip install pytesseract Pillow  "
                "(also requires tesseract-ocr system package)", path,
            )
            return []
        except Exception as exc:
            logger.warning("OCR failed for %s: %s", path, exc)
            return []
        if not text.strip():
            return []
        return [SourceChunk(
            text=text.strip(), source=path, chunk_id=start_id,
            modality="image", metadata={"extraction": "ocr"},
        )]

    _whisper_model: Any = None

    def _ingest_audio(self, path: str, start_id: int = 0) -> list[SourceChunk]:
        """Transcribe audio file. Tries whisper, then skip."""
        try:
            import whisper  # type: ignore[import-untyped]

            if SourceIngester._whisper_model is None:
                SourceIngester._whisper_model = whisper.load_model("base")
            result = SourceIngester._whisper_model.transcribe(path)
            text = result.get("text", "")
        except ImportError:
            logger.warning(
                "Skipping audio %s: pip install openai-whisper", path,
            )
            return []
        except Exception as exc:
            logger.warning("Transcription failed for %s: %s", path, exc)
            return []
        if not text.strip():
            return []
        chunks = self._chunk_text(text.strip(), origin=path, start_id=start_id)
        for c in chunks:
            c.modality = "audio"
            c.metadata = {"extraction": "whisper"}
        return chunks

    def _ingest_video(self, path: str, start_id: int = 0) -> list[SourceChunk]:
        """Extract audio track from video via ffmpeg, then transcribe."""
        import shutil
        import subprocess
        import tempfile

        if not shutil.which("ffmpeg"):
            logger.warning(
                "Skipping video %s: ffmpeg not found in PATH", path,
            )
            return []
        tmpdir = tempfile.mkdtemp(prefix="carl_video_")
        audio_path = Path(tmpdir) / "audio.wav"
        try:
            result = subprocess.run(
                ["ffmpeg", "-i", path, "-vn", "-acodec", "pcm_s16le",
                 "-ar", "16000", "-ac", "1", str(audio_path)],
                capture_output=True, timeout=120,
            )
            if result.returncode != 0 or not audio_path.is_file():
                logger.warning("ffmpeg audio extraction failed for %s", path)
                return []
            chunks = self._ingest_audio(str(audio_path), start_id=start_id)
            for c in chunks:
                c.source = path
                c.modality = "video"
                c.metadata = {"extraction": "ffmpeg+whisper"}
            return chunks
        except Exception as exc:
            logger.warning("Video processing failed for %s: %s", path, exc)
            return []
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    # ------------------------------------------------------------------
    # Remote sources
    # ------------------------------------------------------------------

    def _ingest_github(self, url: str) -> list[SourceChunk]:
        """Clone a GitHub repo to temp dir and ingest."""
        from urllib.parse import urlparse
        parsed = urlparse(url)
        if parsed.scheme != "https" or parsed.netloc not in ("github.com", "www.github.com"):
            raise ValueError(f"Only https://github.com URLs supported: {url}")
        import shutil
        import subprocess
        import tempfile

        tmpdir = tempfile.mkdtemp(prefix="carl_github_")
        try:
            result = subprocess.run(
                ["git", "clone", "--depth", "1", "--", url, tmpdir],
                capture_output=True,
                timeout=60,
            )
            if result.returncode != 0:
                raise ValueError(
                    f"Failed to clone {url}: {result.stderr.decode('utf-8', errors='replace')}"
                )
            return self._ingest_directory(tmpdir)
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    def _ingest_archive(self, path: str) -> list[SourceChunk]:
        """Extract archive and ingest."""
        import shutil
        import tempfile

        if not Path(path).is_file():
            raise FileNotFoundError(f"Archive not found: {path}")
        tmpdir = tempfile.mkdtemp(prefix="carl_archive_")
        try:
            shutil.unpack_archive(path, tmpdir)
            return self._ingest_directory(tmpdir)
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    def _ingest_claw(self, source: str) -> list[SourceChunk]:
        """Ingest OpenClaw dataset. claw://dataset-name

        Downloads from HF and converts to chunks.
        """
        dataset_name = source.removeprefix("claw://")
        hf_id = f"openclaw/{dataset_name}" if "/" not in dataset_name else dataset_name
        return self._ingest_hub_dataset(hf_id)

    # ------------------------------------------------------------------
    # Chunking
    # ------------------------------------------------------------------

    def _chunk_text(
        self, text: str, origin: str, start_id: int = 0
    ) -> list[SourceChunk]:
        """Split text into chunks using content-aware boundaries.

        Strategy selection:
          - .py files: split on top-level class/function definitions
          - .md/.rst files: split on headings
          - Everything else: split on double-newline (paragraph) boundaries
        """
        if origin.endswith(".py"):
            segments = self._split_python(text)
        elif origin.endswith((".md", ".rst")):
            segments = self._split_markdown(text)
        else:
            segments = self._split_paragraphs(text)

        # Merge small segments, split oversized ones
        merged = self._normalize_segments(segments)

        return [
            SourceChunk(text=seg.strip(), source=origin, chunk_id=start_id + i)
            for i, seg in enumerate(merged)
            if seg.strip()
        ]

    @staticmethod
    def _split_python(text: str) -> list[str]:
        """Split Python source on top-level class/function boundaries."""
        pattern = re.compile(r"^(?=(?:class |def |async def ))", re.MULTILINE)
        parts = pattern.split(text)
        return [p for p in parts if p.strip()]

    @staticmethod
    def _split_markdown(text: str) -> list[str]:
        """Split markdown on heading boundaries (# through ####)."""
        pattern = re.compile(r"^(?=#{1,4}\s)", re.MULTILINE)
        parts = pattern.split(text)
        return [p for p in parts if p.strip()]

    @staticmethod
    def _split_paragraphs(text: str) -> list[str]:
        """Split on double-newline boundaries."""
        parts = re.split(r"\n\s*\n", text)
        return [p for p in parts if p.strip()]

    @staticmethod
    def _normalize_segments(segments: list[str]) -> list[str]:
        """Merge undersized segments, split oversized ones."""
        result: list[str] = []
        buffer = ""

        for seg in segments:
            candidate = (buffer + "\n\n" + seg).strip() if buffer else seg.strip()

            if len(candidate) > _MAX_CHUNK_CHARS:
                # Flush buffer first
                if buffer.strip():
                    result.append(buffer.strip())
                    buffer = ""
                # Force-split oversized segment at sentence boundaries
                sentences = re.split(r"(?<=[.!?])\s+", seg)
                sub_buf = ""
                for sent in sentences:
                    if len(sub_buf) + len(sent) > _TARGET_CHUNK_CHARS and sub_buf:
                        result.append(sub_buf.strip())
                        sub_buf = sent
                    else:
                        sub_buf = (sub_buf + " " + sent).strip() if sub_buf else sent
                if sub_buf.strip():
                    buffer = sub_buf.strip()
            elif len(candidate) >= _TARGET_CHUNK_CHARS:
                result.append(candidate)
                buffer = ""
            else:
                buffer = candidate

        if buffer.strip():
            result.append(buffer.strip())

        return result
